from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Helper function to add a title slide
def add_title_slide(title_text, subtitle_text, footer_text):
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    
    # Set subtitle text and style
    tf = subtitle.text_frame
    tf.text = subtitle_text + "\n\n" + footer_text
    
    return slide

# Helper function to add a content slide
def add_content_slide(title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    content = slide.placeholders[1]
    content.text = content_text
    
    return slide

# Slide 1: Title Slide
add_title_slide(
    "CyberShield AI",
    "Real-Time DDoS Threat Intelligence & Monitoring Platform",
    "Presented by: M. Gopi Chandu\nCourse / Department Name\nCollege Name\nYear: 2026"
)

# Slide 2: Problem Statement
add_content_slide(
    "Problem Statement",
    "• Distributed Denial of Service (DDoS) attacks are increasing in scale and complexity.\n"
    "• Organizations lack real-time, visual tools to monitor incoming threat vectors globally.\n"
    "• Existing tools are often too complex, hard to interpret, or lack live visualization.\n"
    "• Delayed threat detection leads to severe network downtime and financial loss."
)

# Slide 3: Objective
add_content_slide(
    "Objective",
    "• Visualize: Build an interactive 3D globe to map attack origins and targets in real-time.\n"
    "• Analyze: Implement an AI-driven threat scoring system to categorize attack severity (Low to Critical).\n"
    "• Monitor: Provide a live stream of incident logs via WebSockets for zero-latency monitoring.\n"
    "• Simplify: Create a seamless, single-pane-of-glass dashboard for Security Operations Center (SOC) analysts."
)

# Slide 4: Key Features
add_content_slide(
    "Key Features",
    "• Interactive 3D Threat Map: Photorealistic globe with attack vectors and impact rings.\n"
    "• Real-Time Incident Feed: Live updating table of incoming attacks.\n"
    "• Automated Threat Categorization: Dynamic severity assignment based on packet rate and protocol.\n"
    "• Analytics Dashboard: Graphical breakdown of attack protocols and system mitigation rates.\n"
    "• Secure Authentication: Role-based access for SOC analysts."
)

# Slide 5: Application Screens
add_content_slide(
    "Application Showcase",
    "(Insert Screenshots Here)\n\n"
    "• Landing Page\n"
    "• Secure Login Portal\n"
    "• 3D Globe Dashboard (Main Command Center)\n"
    "• Dedicated Statistics Page"
)

# Slide 6: Technology Stack
add_content_slide(
    "Technology Stack",
    "Frontend:\n"
    "• React 18, TypeScript, Vite, Tailwind CSS, globe.gl (Three.js), Framer Motion\n\n"
    "Backend:\n"
    "• Python 3, FastAPI, WebSockets, Uvicorn, Pydantic\n\n"
    "Deployment:\n"
    "• GitHub Pages (Client)\n"
    "• Render (API & WebSocket Server)"
)

# Slide 7: System Architecture
add_content_slide(
    "System Architecture",
    "[Client/Browser] <--(WebSockets)--> [FastAPI Backend]\n"
    "[FastAPI Backend] <--(REST API)--> [Authentication Module]\n"
    "[FastAPI Backend] <--(Threat Simulator)--> [Data Generator]\n\n"
    "• Client Tier: React SPA rendering 3D graphics and listening to WebSocket events.\n"
    "• API Tier: FastAPI handling JWT authentication and REST endpoints.\n"
    "• Stream Tier: Continuous WebSocket broadcasting of simulated intrusion events."
)

# Slide 8: Development Methodology
add_content_slide(
    "Development Methodology",
    "1. Requirement Analysis: Identifying SOC needs and visualization requirements.\n"
    "2. Backend Architecture: Developing REST APIs and the WebSocket simulation engine.\n"
    "3. Frontend Integration: Building the React UI and integrating Three.js for the globe.\n"
    "4. Testing & Deployment: Testing latency, ensuring responsiveness, and cloud deployment."
)

# Slide 9: Results Achieved
add_content_slide(
    "Results Achieved",
    "• Successfully simulated and visualized up to 100 concurrent attacks per second.\n"
    "• Sub-second latency between backend generation and frontend rendering.\n"
    "• Fully responsive design functioning across all modern browsers.\n"
    "• Real-time aggregation of statistics (Threat Index, Mitigation Rate)."
)

# Slide 10: Challenges & Solutions
add_content_slide(
    "Challenges & Solutions",
    "Challenge: Managing high-frequency WebSocket updates without freezing the browser.\n"
    "• Solution: Implemented array slicing and efficient React state management to cap rendering limits.\n\n"
    "Challenge: Rendering a photorealistic 3D globe efficiently.\n"
    "• Solution: Utilized WebGL (via globe.gl) to offload rendering to the GPU."
)

# Slide 11: Future Enhancements
add_content_slide(
    "Future Enhancements",
    "• Integration with real IDS/IPS systems (e.g., Snort, Suricata).\n"
    "• Machine Learning integration for predictive attack analysis.\n"
    "• Exportable PDF reports for compliance and auditing.\n"
    "• Multi-tenant architecture for Managed Security Service Providers (MSSPs)."
)

# Slide 12: Thank You
add_content_slide(
    "Thank You!",
    "Any Questions?\n\n"
    "Contact: [Your Email]\n"
    "GitHub: [Your GitHub Profile]\n"
    "Live Demo: [Your App URL]"
)

# Save the presentation
prs.save("CyberShield_AI_Presentation.pptx")
print("Presentation generated successfully!")
